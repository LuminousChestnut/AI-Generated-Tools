from pptx import Presentation

def extract_text_from_pptx(file_path):
    # 打开PowerPoint文件
    presentation = Presentation(file_path)
    slides_text = []

    # 遍历每张幻灯片
    for slide in presentation.slides:
        slide_text = []
        # 遍历幻灯片中的每个形状
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
            elif shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        slide_text.append(run.text)
        slides_text.append(slide_text)
    
    return slides_text

def print_slides_text(slides_text):
    for i, slide_text in enumerate(slides_text):
        print(f"Slide {i + 1}:")
        for text in slide_text:
            print(f"- {text}")
        print("\n")

if __name__ == "__main__":
    file_path = "example.pptx"  # 将此替换为你的PPT文件路径
    slides_text = extract_text_from_pptx(file_path)
    print_slides_text(slides_text)
