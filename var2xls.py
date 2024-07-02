'''
pip install scipy pandas numpy pyreadstat tabula-py python-pptx openpyxl sqlalchemy lxml html5lib
'''

import os
import pandas as pd
import scipy.io
import pyreadstat
import tabula
from pptx import Presentation
from sqlalchemy import create_engine
from io import StringIO, BytesIO
import csv
import sqlite3

def convert_m_to_xls(m_file, output_file):
    data = scipy.io.loadmat(m_file)
    with pd.ExcelWriter(output_file) as writer:
        for key, value in data.items():
            if isinstance(value, np.ndarray):
                df = pd.DataFrame(value)
                df.to_excel(writer, sheet_name=key)

def convert_dta_to_xls(dta_file, output_file):
    df, meta = pyreadstat.read_dta(dta_file)
    df.to_excel(output_file, index=False)

def convert_pdf_to_xls(pdf_file, output_file):
    dfs = tabula.read_pdf(pdf_file, pages='all', multiple_tables=True)
    with pd.ExcelWriter(output_file) as writer:
        for i, df in enumerate(dfs):
            df.to_excel(writer, sheet_name=f'Table_{i+1}')

def convert_ppt_to_xls(ppt_file, output_file):
    presentation = Presentation(ppt_file)
    tables = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                table_data = [[cell.text for cell in row.cells] for row in table.rows]
                df = pd.DataFrame(table_data)
                tables.append(df)
    with pd.ExcelWriter(output_file) as writer:
        for i, df in enumerate(tables):
            df.to_excel(writer, sheet_name=f'Table_{i+1}')

def convert_txt_to_xls(txt_file, output_file, delimiter):
    df = pd.read_csv(txt_file, delimiter=delimiter)
    df.to_excel(output_file, index=False)

def convert_csv_to_xls(csv_file, output_file):
    df = pd.read_csv(csv_file)
    df.to_excel(output_file, index=False)

def convert_sql_to_xls(sql_file, output_file, table_name):
    conn = sqlite3.connect(':memory:')
    with open(sql_file, 'r') as file:
        sql = file.read()
    conn.executescript(sql)
    df = pd.read_sql_query(f"SELECT * FROM {table_name}", conn)
    df.to_excel(output_file, index=False)
    conn.close()

def convert_html_to_xls(html_file, output_file):
    dfs = pd.read_html(html_file)
    with pd.ExcelWriter(output_file) as writer:
        for i, df in enumerate(dfs):
            df.to_excel(writer, sheet_name=f'Table_{i+1}')

def convert_files_to_xls(files, output_dir, delimiter=',', sql_table_name='table'):
    for file in files:
        file_name, file_extension = os.path.splitext(file)
        output_file = os.path.join(output_dir, f"{os.path.basename(file_name)}.xlsx")
        if file_extension == '.m':
            convert_m_to_xls(file, output_file)
        elif file_extension == '.dta':
            convert_dta_to_xls(file, output_file)
        elif file_extension == '.pdf':
            convert_pdf_to_xls(file, output_file)
        elif file_extension == '.pptx':
            convert_ppt_to_xls(file, output_file)
        elif file_extension == '.txt':
            convert_txt_to_xls(file, output_file, delimiter)
        elif file_extension == '.csv':
            convert_csv_to_xls(file, output_file)
        elif file_extension == '.sql':
            convert_sql_to_xls(file, output_file, sql_table_name)
        elif file_extension == '.html':
            convert_html_to_xls(file, output_file)
        else:
            print(f"Unsupported file format: {file_extension}")

if __name__ == "__main__":
    files = ['example.m', 'example.dta', 'example.pdf', 'example.pptx', 'example.txt', 'example.csv', 'example.sql', 'example.html']
    output_dir = 'output_excel_files'
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    convert_files_to_xls(files, output_dir)
