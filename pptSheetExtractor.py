from pptx import Presentation

def extract_tables_from_pptx(file_path):
    # 打开PowerPoint文件
    presentation = Presentation(file_path)
    tables = []

    # 遍历每张幻灯片
    for slide in presentation.slides:
        # 遍历幻灯片中的每个形状
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                table_data = []
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text)
                    table_data.append(row_data)
                tables.append(table_data)
    
    return tables

def print_tables(tables):
    for i, table in enumerate(tables):
        print(f"Table {i + 1}:")
        for row in table:
            print("\t".join(row))
        print("\n")

if __name__ == "__main__":
    file_path = "example.pptx"  # 将此替换为你的PPT文件路径
    tables = extract_tables_from_pptx(file_path)
    print_tables(tables)
