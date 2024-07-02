import os
from pptx import Presentation
from PIL import Image
from io import BytesIO

def save_image(image, path):
    image.save(path)

def extract_images_from_pptx(file_path, separate_folders=True):
    # 打开PowerPoint文件
    presentation = Presentation(file_path)

    # 创建输出文件夹
    base_folder = "ppt_images"
    if not os.path.exists(base_folder):
        os.makedirs(base_folder)

    # 遍历每张幻灯片
    for slide_number, slide in enumerate(presentation.slides, start=1):
        if separate_folders:
            # 为每张幻灯片创建单独的文件夹
            slide_folder = os.path.join(base_folder, f"slide_{slide_number}")
            if not os.path.exists(slide_folder):
                os.makedirs(slide_folder)
        else:
            # 所有图片放在同一个文件夹中
            slide_folder = base_folder

        # 遍历幻灯片中的每个形状
        for shape_number, shape in enumerate(slide.shapes, start=1):
            if hasattr(shape, "image"):
                image = shape.image
                image_bytes = image.blob
                image_stream = BytesIO(image_bytes)
                img = Image.open(image_stream)

                # 保存图片
                image_path = os.path.join(slide_folder, f"image_{slide_number}_{shape_number}.{img.format.lower()}")
                save_image(img, image_path)

if __name__ == "__main__":
    file_path = "example.pptx"  # 将此替换为你的PPT文件路径
    separate_folders = True  # 设置为False以将所有图片放在同一个文件夹中
    extract_images_from_pptx(file_path, separate_folders)
